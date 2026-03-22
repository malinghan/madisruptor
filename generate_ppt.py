#!/usr/bin/env python3
"""
Generate a comprehensive PPTX presentation:
"Java 高性能队列 Disruptor 原理详解"
Based on the MaDisruptor GitHub tutorial.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─────────── Colour Palette (dark-tech style) ───────────
BG_DARK      = RGBColor(0x1B, 0x1B, 0x2F)   # deep navy
BG_CARD      = RGBColor(0x24, 0x24, 0x3E)   # card background
ACCENT       = RGBColor(0x00, 0xD2, 0xFF)   # cyan accent
ACCENT2      = RGBColor(0x7C, 0x4D, 0xFF)   # purple accent
ACCENT3      = RGBColor(0xFF, 0x6B, 0x6B)   # coral / red
ACCENT4      = RGBColor(0x51, 0xCF, 0x66)   # green
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xBB, 0xBB, 0xCC)
CODE_BG      = RGBColor(0x18, 0x18, 0x28)
YELLOW       = RGBColor(0xFF, 0xD4, 0x3B)
ORANGE       = RGBColor(0xFF, 0x92, 0x2B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ═══════════════════════════════════════════════════════════
#  Helper functions
# ═══════════════════════════════════════════════════════════

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_text(slide, left, top, width, height, lines, font_size=16, color=WHITE, line_spacing=1.3, font_name="Microsoft YaHei"):
    """lines: list of (text, color, bold, font_size_override)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, clr, bld, fsz = item, color, False, font_size
        elif len(item) == 2:
            txt, clr = item; bld = False; fsz = font_size
        elif len(item) == 3:
            txt, clr, bld = item; fsz = font_size
        else:
            txt, clr, bld, fsz = item
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(fsz)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1) * 2)
    return txBox

def add_code_block(slide, left, top, width, height, code_text, font_size=11):
    """Add a styled code block."""
    rect = add_rect(slide, left, top, width, height, CODE_BG, ACCENT, Pt(1))
    rect.text_frame.word_wrap = True
    rect.text_frame.margin_left = Inches(0.2)
    rect.text_frame.margin_right = Inches(0.2)
    rect.text_frame.margin_top = Inches(0.15)
    rect.text_frame.margin_bottom = Inches(0.15)
    for i, line in enumerate(code_text.strip().split('\n')):
        if i == 0:
            p = rect.text_frame.paragraphs[0]
        else:
            p = rect.text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = LIGHT_GRAY
        p.font.name = "Consolas"
        p.space_after = Pt(1)
    return rect

def add_chapter_header(slide, chapter_num, chapter_title, subtitle=""):
    """Add a consistent chapter header bar."""
    # Top accent line
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT)
    # Chapter number badge
    badge = add_rect(slide, Inches(0.6), Inches(0.35), Inches(0.7), Inches(0.5), ACCENT)
    badge.text_frame.paragraphs[0].text = f"{chapter_num:02d}"
    badge.text_frame.paragraphs[0].font.size = Pt(22)
    badge.text_frame.paragraphs[0].font.color.rgb = BG_DARK
    badge.text_frame.paragraphs[0].font.bold = True
    badge.text_frame.paragraphs[0].font.name = "Consolas"
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    badge.text_frame.paragraphs[0].space_before = Pt(4)
    # Title
    add_textbox(slide, Inches(1.5), Inches(0.3), Inches(10), Inches(0.6),
                chapter_title, font_size=30, color=WHITE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(1.5), Inches(0.82), Inches(10), Inches(0.4),
                    subtitle, font_size=15, color=LIGHT_GRAY)

def add_section_title(slide, title_text):
    """Full-page section divider slide."""
    set_slide_bg(slide, BG_DARK)
    # Center accent bar
    add_rect(slide, Inches(5.5), Inches(2.8), Inches(2.3), Inches(0.04), ACCENT)
    add_textbox(slide, Inches(1), Inches(3.1), Inches(11.3), Inches(1.2),
                title_text, font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

def new_slide():
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide)
    return slide

def add_bullet_card(slide, left, top, width, height, title, items, title_color=ACCENT, item_color=WHITE, font_size=15, title_size=18):
    """Card with title and bullet items."""
    card = add_rect(slide, left, top, width, height, BG_CARD, RGBColor(0x3A,0x3A,0x55), Pt(1))
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.1)
    # Title
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.space_after = Pt(8)
    # Items
    for item in items:
        p = tf.add_paragraph()
        if isinstance(item, tuple):
            p.text = item[0]
            p.font.color.rgb = item[1]
        else:
            p.text = f"  {item}"
            p.font.color.rgb = item_color
        p.font.size = Pt(font_size)
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(4)
    return card

def add_icon_text(slide, left, top, icon_char, text, icon_color=ACCENT, text_color=WHITE, font_size=18):
    """Simple icon + text combo."""
    add_textbox(slide, left, top, Inches(0.4), Inches(0.4), icon_char, font_size=font_size+4, color=icon_color, bold=True, font_name="Segoe UI Emoji")
    add_textbox(slide, left + Inches(0.45), top + Inches(0.02), Inches(5), Inches(0.4), text, font_size=font_size, color=text_color, bold=False)


# ╔═══════════════════════════════════════════════════════════════╗
# ║                    SLIDE 1: COVER                            ║
# ╚═══════════════════════════════════════════════════════════════╝

slide = new_slide()
# Big gradient-like band
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT)
add_rect(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), ACCENT2)

# Decorative circles
for cx, cy, r, c in [(Inches(10.5), Inches(1.2), Inches(2.5), ACCENT), (Inches(11.5), Inches(2.8), Inches(1.5), ACCENT2)]:
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, r, r)
    shape.fill.solid()
    shape.fill.fore_color.rgb = c
    shape.fill.fore_color.brightness = 0.85
    shape.line.fill.background()

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(9), Inches(0.7),
            "JAVA 高性能队列", font_size=22, color=ACCENT, bold=True)
add_textbox(slide, Inches(0.8), Inches(2.2), Inches(10), Inches(1.5),
            "Disruptor 原理详解", font_size=52, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(3.8), Inches(10), Inches(0.8),
            "从原理到实战 · 从入门到精通", font_size=24, color=LIGHT_GRAY)

# Bottom info bar
add_rect(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.04), RGBColor(0x3A,0x3A,0x55))
lines = [
    ("GitHub: malinghan/madisruptor", LIGHT_GRAY),
    ("基于 LMAX Disruptor 框架 · 8 章文档 + 8 个 Demo · Java 11+", LIGHT_GRAY),
]
add_multiline_text(slide, Inches(0.8), Inches(5.7), Inches(11), Inches(1), lines, font_size=14)


# ╔═══════════════════════════════════════════════════════════════╗
# ║                  SLIDE 2: AGENDA / TOC                       ║
# ╚═══════════════════════════════════════════════════════════════╝

slide = new_slide()
add_chapter_header(slide, 0, "目录 · AGENDA", "8 个章节，循序渐进掌握 Disruptor")

chapters = [
    ("01", "为什么需要 Disruptor",    "传统队列痛点分析"),
    ("02", "核心概念与架构",          "9 大组件详解"),
    ("03", "RingBuffer 原理详解",     "位运算 · 缓存行填充 · 预分配"),
    ("04", "Sequence 与并发控制",     "CAS 无锁 · 内存屏障"),
    ("05", "WaitStrategy 等待策略",   "5 种策略对比与选型"),
    ("06", "EventProcessor 事件处理", "批量消费 · 管道 · 菱形依赖"),
    ("07", "性能测试与最佳实践",      "Benchmark · 调优 · 常见陷阱"),
    ("08", "手写简易 Disruptor",      "从零实现核心功能"),
]
x_start = Inches(0.6)
y_start = Inches(1.5)
card_w = Inches(3.0)
card_h = Inches(0.7)
gap = Inches(0.12)

for i, (num, title, desc) in enumerate(chapters):
    col = i % 2
    row = i // 2
    x = x_start + col * (card_w + Inches(0.2))
    y = y_start + row * (card_h + gap)
    
    card = add_rect(slide, x, y, card_w, card_h, BG_CARD, RGBColor(0x3A,0x3A,0x55), Pt(1))
    # Number badge
    badge = add_rect(slide, x + Inches(0.1), y + Inches(0.12), Inches(0.45), Inches(0.45), ACCENT if col == 0 else ACCENT2)
    badge.text_frame.paragraphs[0].text = num
    badge.text_frame.paragraphs[0].font.size = Pt(16)
    badge.text_frame.paragraphs[0].font.color.rgb = BG_DARK
    badge.text_frame.paragraphs[0].font.bold = True
    badge.text_frame.paragraphs[0].font.name = "Consolas"
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_textbox(slide, x + Inches(0.65), y + Inches(0.06), Inches(2.2), Inches(0.35),
                title, font_size=13, color=WHITE, bold=True)
    add_textbox(slide, x + Inches(0.65), y + Inches(0.38), Inches(2.2), Inches(0.3),
                desc, font_size=10, color=LIGHT_GRAY)

# Learning path on right side
path_x = Inches(6.8)
add_textbox(slide, path_x, Inches(1.5), Inches(5.5), Inches(0.4), "学习路线", font_size=20, color=ACCENT, bold=True)

stages = [
    ("第一阶段 · 快速入门", "1-2 小时", "Ch01-02 + Demo01", ACCENT4),
    ("第二阶段 · 原理深入", "3-5 小时", "Ch03-06 + Demo02-05", ACCENT),
    ("第三阶段 · 实战精通", "3-5 小时", "Ch07-08 + Demo06-08", ACCENT2),
]
for i, (stage, time, content, color) in enumerate(stages):
    sy = Inches(2.1) + i * Inches(1.25)
    add_rect(slide, path_x, sy, Inches(5.5), Inches(1.0), BG_CARD, color, Pt(1.5))
    add_textbox(slide, path_x + Inches(0.2), sy + Inches(0.1), Inches(4), Inches(0.35), stage, font_size=15, color=color, bold=True)
    add_textbox(slide, path_x + Inches(0.2), sy + Inches(0.45), Inches(4), Inches(0.25), content, font_size=12, color=LIGHT_GRAY)
    add_textbox(slide, path_x + Inches(4.2), sy + Inches(0.1), Inches(1.1), Inches(0.35), time, font_size=12, color=YELLOW, bold=True, alignment=PP_ALIGN.RIGHT)
    # connector arrow
    if i < 2:
        add_textbox(slide, path_x + Inches(2.5), sy + Inches(1.0), Inches(0.5), Inches(0.25), "▼", font_size=14, color=color, alignment=PP_ALIGN.CENTER)

# Bottom note
add_rect(slide, Inches(6.8), Inches(5.95), Inches(5.5), Inches(0.8), BG_CARD, ORANGE, Pt(1))
add_textbox(slide, Inches(7.0), Inches(6.0), Inches(5), Inches(0.7), "配套 Demo: 每章配有可直接运行的 Java Demo 代码\nmvn exec:java -Dexec.mainClass=\"com.madisruptor.demo0X...\"", font_size=11, color=LIGHT_GRAY)


# ╔═══════════════════════════════════════════════════════════════╗
# ║           CHAPTER 1: 为什么需要 Disruptor                     ║
# ╚═══════════════════════════════════════════════════════════════╝

# --- 1-1: Section divider ---
slide = new_slide()
add_section_title(slide, "01  为什么需要 Disruptor")

# --- 1-2: 传统队列痛点 ---
slide = new_slide()
add_chapter_header(slide, 1, "传统并发队列的四大痛点")

problems = [
    ("锁竞争 Lock Contention", "所有生产者/消费者竞争同一把锁\n频繁的线程阻塞唤醒，上下文切换开销巨大", ACCENT3),
    ("伪共享 False Sharing", "putIndex 和 takeIndex 在同一缓存行\n一个线程修改 → 整个 Cache Line 失效", ORANGE),
    ("频繁 GC", "LinkedBlockingQueue 每次入队 new Node\n大量临时对象 → GC 压力 → STW 暂停", YELLOW),
    ("缓存不友好", "链表节点内存随机分布\nCPU 预取失效，频繁 Cache Miss", ACCENT2),
]
for i, (title, desc, color) in enumerate(problems):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + col * Inches(6.2)
    y = Inches(1.4) + row * Inches(2.6)
    card = add_rect(slide, x, y, Inches(5.9), Inches(2.3), BG_CARD, color, Pt(2))
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.2)
    # Number
    num_shape = add_rect(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.45), Inches(0.45), color)
    num_shape.text_frame.paragraphs[0].text = f"0{i+1}"
    num_shape.text_frame.paragraphs[0].font.size = Pt(18)
    num_shape.text_frame.paragraphs[0].font.color.rgb = BG_DARK
    num_shape.text_frame.paragraphs[0].font.bold = True
    num_shape.text_frame.paragraphs[0].font.name = "Consolas"
    num_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_textbox(slide, x + Inches(0.75), y + Inches(0.15), Inches(4.8), Inches(0.4), title, font_size=18, color=WHITE, bold=True)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.7), Inches(5.2), Inches(1.4), desc, font_size=14, color=LIGHT_GRAY)

# --- 1-3: 伪共享图解 ---
slide = new_slide()
add_chapter_header(slide, 1, "伪共享 (False Sharing) 图解")

add_code_block(slide, Inches(0.5), Inches(1.4), Inches(6), Inches(2.0),
"""CPU Cache Line (64 bytes)
┌──────────────────────────────────────┐
│  head (8B)  │  tail (8B)  │  其他... │
└──────────────────────────────────────┘
       ↑               ↑
    Thread A         Thread B
    (生产者)          (消费者)

Thread A 修改 head → 整个缓存行失效
→ Thread B 必须重新加载 tail (虽然没被改！)""", font_size=13)

add_textbox(slide, Inches(7), Inches(1.4), Inches(5.5), Inches(0.4), "内存布局对比", font_size=18, color=ACCENT, bold=True)

add_code_block(slide, Inches(7), Inches(1.9), Inches(5.8), Inches(1.5),
"""LinkedBlockingQueue (链表):
[Node1] → ... → [Node2] → ... → [Node3]
 0x100          0x900          0x400
 ↑ 随机分布，缓存不友好

RingBuffer (数组):
[Slot0][Slot1][Slot2][Slot3][Slot4][Slot5]
 0x100 0x120  0x140  0x160  0x180  0x1A0
 ↑ 连续分布，CPU 缓存友好""", font_size=12)

# Solution table
add_textbox(slide, Inches(0.5), Inches(3.7), Inches(12), Inches(0.4), "Disruptor 的解决方案", font_size=20, color=ACCENT, bold=True)

table_data = [
    ("传统队列问题", "Disruptor 解决方案"),
    ("锁竞争", "CAS 无锁 + Sequence 协调"),
    ("伪共享", "Cache Line Padding (缓存行填充)"),
    ("频繁 GC", "预分配 Event 对象，循环复用"),
    ("缓存不友好", "数组 + 顺序访问"),
    ("固定等待", "可插拔 WaitStrategy"),
    ("逐条处理", "批量消费 (Batching)"),
]
table_shape = slide.shapes.add_table(len(table_data), 2, Inches(0.5), Inches(4.2), Inches(12), Inches(2.8))
table = table_shape.table
table.columns[0].width = Inches(4)
table.columns[1].width = Inches(8)
for r, (c1, c2) in enumerate(table_data):
    for ci, val in enumerate([c1, c2]):
        cell = table.cell(r, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.name = "Microsoft YaHei"
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = BG_DARK
            else:
                p.font.color.rgb = WHITE if ci == 0 else ACCENT4
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if r == 0 else (BG_CARD if r % 2 == 1 else RGBColor(0x2A,0x2A,0x44))


# ╔═══════════════════════════════════════════════════════════════╗
# ║           CHAPTER 2: 核心概念与架构                            ║
# ╚═══════════════════════════════════════════════════════════════╝

slide = new_slide()
add_section_title(slide, "02  核心概念与架构")

# --- 2-1: Architecture diagram ---
slide = new_slide()
add_chapter_header(slide, 2, "Disruptor 核心架构", "生产者 → Sequencer → RingBuffer → SequenceBarrier → 消费者")

add_code_block(slide, Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.8),
"""                              Disruptor 核心架构

    Producer(s)
    ┌──────────┐  ┌──────────┐
    │Producer 1│  │Producer 2│  ...
    └────┬─────┘  └────┬─────┘
         │              │
         ▼              ▼
    ┌─────────────────────────────────┐
    │          Sequencer              │  协调生产者发布
    │   SingleProducerSequencer       │  单生产者: 无竞争, 最高性能
    │   MultiProducerSequencer        │  多生产者: CAS 无锁竞争
    └─────────────┬───────────────────┘
                  │
                  ▼
    ┌─────────────────────────────────┐
    │         RingBuffer              │  预分配的环形数组
    │    ┌───┬───┬───┬───┬───┬───┐   │  size = 2^n (位运算取模)
    │    │ 0 │ 1 │ 2 │ 3 │ 4 │ 5 │   │
    │    └───┴───┴───┴───┴───┴───┘   │
    └─────────────┬───────────────────┘
                  │
                  ▼
    ┌─────────────────────────────────┐
    │      SequenceBarrier            │  协调消费者等待
    │      + WaitStrategy             │  BusySpin / Yield / Block / ...
    └─────────────┬───────────────────┘
         ┌────────┼────────┐
         ▼        ▼        ▼
    ┌──────────┐┌──────────┐┌──────────┐
    │Consumer 1││Consumer 2││Consumer 3│
    │(EventHdl)││(EventHdl)││(WorkHdl) │
    └──────────┘└──────────┘└──────────┘""", font_size=11)


# --- 2-2: 核心概念速查表 ---
slide = new_slide()
add_chapter_header(slide, 2, "核心概念速查表", "9 大组件一览")

concepts = [
    ("RingBuffer",      "环形数组，存储 Event 的容器\n大小必须为 2 的 N 次方",             ACCENT),
    ("Event",           "RingBuffer 中的数据载体 (POJO)\n预创建，循环复用，零 GC",          ACCENT4),
    ("EventFactory",    "创建 Event 的工厂\n用于预填充 RingBuffer",                       ACCENT4),
    ("Sequence",        "递增序号，标识 RingBuffer 位置\n带 CacheLine Padding 避免伪共享",  YELLOW),
    ("Sequencer",       "协调生产者的序号分配\nSingle (无锁) / Multi (CAS)",               YELLOW),
    ("SequenceBarrier", "协调消费者等待可用事件\n内部委托给 WaitStrategy",                   ORANGE),
    ("WaitStrategy",    "消费者等待策略\nBusySpin / Yield / Sleep / Block",                ORANGE),
    ("EventHandler",    "事件处理器接口\nonEvent(event, seq, endOfBatch)",                 ACCENT2),
    ("EventProcessor",  "事件处理的执行单元\nBatchEventProcessor / WorkProcessor",          ACCENT2),
]

for i, (name, desc, color) in enumerate(concepts):
    col = i % 3
    row = i // 3
    x = Inches(0.4) + col * Inches(4.2)
    y = Inches(1.4) + row * Inches(1.9)
    card = add_rect(slide, x, y, Inches(4.0), Inches(1.7), BG_CARD, color, Pt(1.5))
    add_textbox(slide, x + Inches(0.2), y + Inches(0.12), Inches(3.6), Inches(0.35), name, font_size=16, color=color, bold=True, font_name="Consolas")
    add_textbox(slide, x + Inches(0.2), y + Inches(0.55), Inches(3.6), Inches(1.0), desc, font_size=12, color=LIGHT_GRAY)


# --- 2-3: 事件发布与消费流程 ---
slide = new_slide()
add_chapter_header(slide, 2, "事件发布与消费流程")

add_textbox(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(0.4), "生产者发布流程", font_size=18, color=ACCENT, bold=True)
add_code_block(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(2.8),
"""// Step 1: 申请序号
long sequence = ringBuffer.next();

// Step 2: 获取该位置的 Event 对象 (预分配)
OrderEvent event = ringBuffer.get(sequence);

// Step 3: 填充数据 (覆盖写入, 不 new 对象!)
event.setOrderId(123L);
event.setPrice(99.9);

// Step 4: 发布事件 (更新 cursor, 通知消费者)
ringBuffer.publish(sequence);""", font_size=12)

# Recommended way
add_textbox(slide, Inches(0.5), Inches(4.8), Inches(5.5), Inches(0.3), "推荐: EventTranslator 方式 (更安全)", font_size=14, color=ACCENT4, bold=True)
add_code_block(slide, Inches(0.5), Inches(5.2), Inches(5.8), Inches(1.6),
"""// 一步完成, 不会忘记 publish
ringBuffer.publishEvent(
    (event, sequence, data) -> {
        event.setOrderId(data);
    },
    orderId
);""", font_size=12)

add_textbox(slide, Inches(6.8), Inches(1.3), Inches(5.5), Inches(0.4), "消费者消费流程", font_size=18, color=ACCENT2, bold=True)
add_code_block(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(2.8),
"""// Step 1: 等待可用序号 (可能批量返回!)
long available = barrier.waitFor(nextSeq);

// Step 2: 批量处理
while (nextSeq <= available) {
    Event e = ringBuffer.get(nextSeq);
    handler.onEvent(e, nextSeq,
        nextSeq == available); // endOfBatch
    nextSeq++;
}

// Step 3: 更新消费进度
sequence.set(available);""", font_size=12)

add_textbox(slide, Inches(6.8), Inches(4.8), Inches(5.5), Inches(0.3), "批量消费的优势", font_size=14, color=ACCENT4, bold=True)
add_code_block(slide, Inches(6.8), Inches(5.2), Inches(5.8), Inches(1.6),
"""// 生产者发了 seq 5,6,7,8,9,10
// 消费者在 seq 4

// 传统: 6次 wait → 6次 process
// Disruptor: 1次 waitFor(5)→返回10
//   → 一次性处理 6 个事件!""", font_size=12)


# ╔═══════════════════════════════════════════════════════════════╗
# ║           CHAPTER 3: RingBuffer 原理详解                      ║
# ╚═══════════════════════════════════════════════════════════════╝

slide = new_slide()
add_section_title(slide, "03  RingBuffer 原理详解")

# --- 3-1: 位运算取模 ---
slide = new_slide()
add_chapter_header(slide, 3, "位运算取模的妙用", "bufferSize 必须是 2 的幂")

add_code_block(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(3.5),
"""// 普通取模 (除法: 20-90 CPU 时钟周期)
index = sequence % bufferSize;

// 位运算取模 (位与: 1 个时钟周期!)
index = sequence & (bufferSize - 1);

// 原理 (bufferSize=8):
// bufferSize   = 8  → 二进制: 1000
// bufferSize-1 = 7  → 二进制: 0111 (掩码)

// sequence=13 → 二进制: 1101
// 13 & 7      → 二进制: 0101 = 5
// 13 % 8      →                 5 ✓

// 性能差距: 20~90 倍!
// 每秒数百万次操作, 差距非常显著""", font_size=13)

# RingBuffer visual
add_textbox(slide, Inches(7), Inches(1.4), Inches(5), Inches(0.4), "环形缓冲区示意图", font_size=18, color=ACCENT, bold=True)
add_code_block(slide, Inches(7), Inches(1.9), Inches(5.5), Inches(3.0),
"""      ┌───┐
 ┌───►│ 7 │◄───┐
 │    └───┘    │
┌───┐        ┌───┐
│ 6 │        │ 0 │  ← cursor
└───┘        └───┘
┌───┐        ┌───┐
│ 5 │        │ 1 │
└───┘        └───┘
 │    ┌───┐    │
 └───►│ 4 │◄───┘
      └───┘
 ┌───┐    ┌───┐
 │ 3 │────│ 2 │
 └───┘    └───┘

 RingBuffer (size=8, 即 2^3)
 index = sequence & 0x07""", font_size=13)

# Key advantages
add_textbox(slide, Inches(0.5), Inches(5.2), Inches(12), Inches(0.4), "RingBuffer 三大优势", font_size=18, color=YELLOW, bold=True)
advantages = [
    ("固定大小，预分配", "一次性分配所有内存，GC 友好，对象在老年代长期存活", ACCENT4),
    ("连续内存，缓存友好", "数组结构 + 顺序访问，CPU 预取高效", ACCENT),
    ("位运算循环", "只维护递增的 sequence，无需移动数据", ACCENT2),
]
for i, (t, d, c) in enumerate(advantages):
    x = Inches(0.5) + i * Inches(4.2)
    card = add_rect(slide, x, Inches(5.7), Inches(3.9), Inches(1.2), BG_CARD, c, Pt(1.5))
    add_textbox(slide, x + Inches(0.15), Inches(5.8), Inches(3.6), Inches(0.3), t, font_size=14, color=c, bold=True)
    add_textbox(slide, x + Inches(0.15), Inches(6.15), Inches(3.6), Inches(0.6), d, font_size=11, color=LIGHT_GRAY)


# --- 3-2: 缓存行填充 ---
slide = new_slide()
add_chapter_header(slide, 3, "缓存行填充 (Cache Line Padding)", "消除伪共享的核心技术")

add_code_block(slide, Inches(0.5), Inches(1.4), Inches(6), Inches(4.0),
"""// Disruptor Sequence 类的缓存行填充设计

abstract class RingBufferPad {
    protected long p1, p2, p3, p4, p5, p6, p7;
    // 前置填充: 7 × 8 = 56 bytes
}

abstract class RingBufferFields<E>
        extends RingBufferPad {
    // 实际字段
    private final long indexMask;
    private final Object[] entries;
    protected final int bufferSize;
    protected final Sequencer sequencer;
}

public final class RingBuffer<E>
        extends RingBufferFields<E> {
    // 后置填充: 7 × 8 = 56 bytes
    protected long p1, p2, p3, p4, p5, p6, p7;
}""", font_size=12)

add_textbox(slide, Inches(7), Inches(1.4), Inches(5.5), Inches(0.4), "填充后的内存布局", font_size=18, color=ACCENT, bold=True)
add_code_block(slide, Inches(7), Inches(1.9), Inches(5.5), Inches(2.0),
"""Cache Line 1     Cache Line 2     Cache Line 3
┌──────────────┐ ┌──────────────┐ ┌──────────────┐
│ p1 p2 p3 p4 │ │              │ │ p5 p6 p7     │
│ p5 p6 p7    │ │  value (8B)  │ │              │
│ (前置填充)    │ │ (独占缓存行!) │ │ (后置填充)    │
└──────────────┘ └──────────────┘ └──────────────┘

→ value 独占一个 64-byte 缓存行
→ 不同线程修改不同 Sequence
→ 互不干扰, 零伪共享!""", font_size=12)

add_textbox(slide, Inches(7), Inches(4.2), Inches(5.5), Inches(0.4), "预分配与对象复用", font_size=18, color=ACCENT4, bold=True)
add_code_block(slide, Inches(7), Inches(4.7), Inches(5.5), Inches(2.5),
"""// 传统方式: 每次 new 对象 (GC 压力大)
queue.put(new OrderEvent(id, price));

// Disruptor: 初始化时一次性创建
for (int i = 0; i < bufferSize; i++) {
    entries[i] = factory.newInstance();
}
// 发布时只是覆盖写入, 不 new 对象!
long seq = ringBuffer.next();
OrderEvent e = ringBuffer.get(seq);
e.setOrderId(id);  // 覆盖
ringBuffer.publish(seq);""", font_size=11)


# ╔═══════════════════════════════════════════════════════════════╗
# ║           CHAPTER 4: Sequence 与并发控制                      ║
# ╚═══════════════════════════════════════════════════════════════╝

slide = new_slide()
add_section_title(slide, "04  Sequence 与并发控制")

slide = new_slide()
add_chapter_header(slide, 4, "Single vs Multi ProducerSequencer", "无锁并发的核心设计")

add_textbox(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.4), "SingleProducerSequencer (最高性能)", font_size=17, color=ACCENT4, bold=True)
add_code_block(slide, Inches(0.5), Inches(1.8), Inches(6), Inches(3.5),
"""// 单生产者: 无竞争, 无CAS, 非volatile读
public long next() {
    long nextValue = this.nextValue;
    // ↑ 非volatile! 性能极高
    long nextSeq = nextValue + 1;
    long wrapPoint = nextSeq - bufferSize;

    // 只在必要时查询消费者最小 sequence
    if (wrapPoint > cachedGatingSequence) {
        long minSeq;
        while (wrapPoint > (minSeq =
                getMinGatingSequence())) {
            LockSupport.parkNanos(1);
        }
        cachedGatingSequence = minSeq;
    }
    this.nextValue = nextSeq;
    return nextSeq;
}""", font_size=11)

add_textbox(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(0.4), "MultiProducerSequencer (CAS竞争)", font_size=17, color=ORANGE, bold=True)
add_code_block(slide, Inches(6.8), Inches(1.8), Inches(6), Inches(3.5),
"""// 多生产者: CAS 无锁竞争
public long next() {
    long current;
    long next;
    do {
        current = cursor.get();
        // ↑ volatile读
        next = current + 1;

        long wrapPoint = next - bufferSize;
        if (wrapPoint > cachedGatingSeq) {
            // 等待消费者...
            continue;
        } else if (cursor.compareAndSet(
                current, next)) {
            break; // CAS成功!
        }
        // CAS失败则重试
    } while (true);
    return next;
}""", font_size=11)

# Bottom comparison
add_rect(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.5), BG_CARD, ACCENT, Pt(1))
add_textbox(slide, Inches(0.7), Inches(5.7), Inches(11.5), Inches(0.3), "关键区别与建议", font_size=15, color=ACCENT, bold=True)
add_textbox(slide, Inches(0.7), Inches(6.05), Inches(11.5), Inches(0.9),
"SingleProducer: nextValue 非 volatile, 无 CAS, 使用 cachedValue 减少 volatile 读 → 性能最高\n"
"MultiProducer: cursor 是 volatile, 使用 CAS 竞争 + availableBuffer 标记可见性 → 有额外开销\n"
"最佳实践: 即使有多个数据来源, 也尽量汇聚到一个线程后再发布 (变成单生产者)", font_size=12, color=LIGHT_GRAY)


# ╔═══════════════════════════════════════════════════════════════╗
# ║           CHAPTER 5: WaitStrategy 等待策略                    ║
# ╚═══════════════════════════════════════════════════════════════╝

slide = new_slide()
add_section_title(slide, "05  WaitStrategy 等待策略")

slide = new_slide()
add_chapter_header(slide, 5, "五种等待策略对比", "消费者没有新事件时怎么等？")

strategies = [
    ("BusySpinWaitStrategy", "纯自旋",        "极高", "极低(ns)",   "超低延迟 + CPU充足",      ACCENT3, "while(...) onSpinWait();"),
    ("YieldingWaitStrategy", "自旋100次+yield","高",   "低(亚μs)",   "低延迟首选, 推荐",        ACCENT4, "spin 100 → Thread.yield()"),
    ("SleepingWaitStrategy", "三段式",         "低",   "中(μs)",     "后台任务如日志",          ACCENT, "spin → yield → sleep(1ns)"),
    ("BlockingWaitStrategy", "Lock+Condition", "最低", "高(ms)",     "CPU受限/吞吐优先",        ACCENT2, "lock → condition.await()"),
    ("TimeoutBlocking",      "超时阻塞",       "低",   "可控",       "需超时处理如心跳",        ORANGE, "condition.await(timeout)"),
]

# Table
table_shape = slide.shapes.add_table(6, 6, Inches(0.4), Inches(1.4), Inches(12.5), Inches(3.7))
table = table_shape.table
headers = ["策略", "原理", "CPU占用", "延迟", "适用场景", "核心实现"]
widths = [Inches(2.5), Inches(1.8), Inches(1.0), Inches(1.2), Inches(2.8), Inches(3.2)]
for ci, w in enumerate(widths):
    table.columns[ci].width = w

for ci, h in enumerate(headers):
    cell = table.cell(0, ci)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = BG_DARK; p.font.name = "Microsoft YaHei"
    cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT

for ri, (name, principle, cpu, latency, usecase, color, impl) in enumerate(strategies):
    vals = [name, principle, cpu, latency, usecase, impl]
    for ci, val in enumerate(vals):
        cell = table.cell(ri+1, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = color if ci == 0 else LIGHT_GRAY
            p.font.name = "Consolas" if ci in (0, 5) else "Microsoft YaHei"
            p.font.bold = (ci == 0)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG_CARD if ri % 2 == 0 else RGBColor(0x2A,0x2A,0x44)

# Decision guide
add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12), Inches(0.4), "选型决策指南", font_size=17, color=YELLOW, bold=True)
add_code_block(slide, Inches(0.5), Inches(5.7), Inches(12), Inches(1.4),
"""延迟 < 1μs?  → CPU充足? → BusySpin    |  CPU受限? → Yielding
延迟 < 100μs? → Yielding 或 Sleeping
延迟 > 100μs? → 需超时? → TimeoutBlocking  |  不需要? → Blocking 或 Sleeping

Log4j2 AsyncLogger 默认使用 TimeoutBlockingWaitStrategy (日志对延迟不敏感, 需节省CPU)""", font_size=12)


# ╔═══════════════════════════════════════════════════════════════╗
# ║           CHAPTER 6: EventProcessor 事件处理                  ║
# ╚═══════════════════════════════════════════════════════════════╝

slide = new_slide()
add_section_title(slide, "06  EventProcessor 事件处理")

# --- 6-1: Consumer patterns ---
slide = new_slide()
add_chapter_header(slide, 6, "消费者组合模式", "Disruptor DSL 灵活构建处理管道")

patterns = [
    ("串行 Pipeline", "[P] → [C1] → [C2] → [C3]", ".handleEventsWith(h1).then(h2).then(h3)", ACCENT),
    ("并行 Broadcast", "[P] → [C1]\n     → [C2]\n     → [C3]", ".handleEventsWith(h1, h2, h3)", ACCENT4),
    ("竞争 WorkerPool", "[P] → [W1] 或 [W2] 或 [W3]\n(每个事件只处理一次)", ".handleEventsWithWorkerPool(w1,w2,w3)", ORANGE),
    ("菱形 Diamond", "     → [C1] ─┐\n[P]            ├→ [C3]\n     → [C2] ─┘", ".handleEventsWith(h1,h2).then(h3)", ACCENT2),
]

for i, (title, diagram, code, color) in enumerate(patterns):
    col = i % 2
    row = i // 2
    x = Inches(0.4) + col * Inches(6.4)
    y = Inches(1.4) + row * Inches(2.9)
    card = add_rect(slide, x, y, Inches(6.1), Inches(2.7), BG_CARD, color, Pt(2))
    add_textbox(slide, x + Inches(0.2), y + Inches(0.1), Inches(5.5), Inches(0.35), title, font_size=16, color=color, bold=True)
    add_code_block(slide, x + Inches(0.15), y + Inches(0.5), Inches(5.8), Inches(1.2), diagram, font_size=12)
    add_textbox(slide, x + Inches(0.2), y + Inches(1.85), Inches(5.6), Inches(0.6), code, font_size=10, color=YELLOW, font_name="Consolas")


# ╔═══════════════════════════════════════════════════════════════╗
# ║           CHAPTER 7: 性能测试与最佳实践                        ║
# ╚═══════════════════════════════════════════════════════════════╝

slide = new_slide()
add_section_title(slide, "07  性能测试与最佳实践")

# --- 7-1: Benchmark results ---
slide = new_slide()
add_chapter_header(slide, 7, "性能基准测试结果", "Demo06: Disruptor vs ArrayBlockingQueue  (100万事件)")

# 1P1C Results
add_textbox(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(0.4), "单生产者-单消费者 (1P-1C)", font_size=18, color=ACCENT, bold=True)
# Bar chart simulation
bar_data_1p1c = [("Disruptor", 11, ACCENT4), ("ArrayBlockingQueue", 128, ACCENT3)]
for i, (name, ms, color) in enumerate(bar_data_1p1c):
    y = Inches(1.9) + i * Inches(1.0)
    add_textbox(slide, Inches(0.5), y, Inches(2.5), Inches(0.35), name, font_size=13, color=WHITE, bold=True)
    bar_width = max(Inches(0.3), Inches(ms / 128.0 * 3.5))
    add_rect(slide, Inches(3.2), y + Inches(0.05), bar_width, Inches(0.35), color)
    add_textbox(slide, Inches(3.2) + bar_width + Inches(0.1), y, Inches(1.5), Inches(0.35), f"{ms} ms", font_size=14, color=color, bold=True)

# Big number
add_rect(slide, Inches(0.5), Inches(4.0), Inches(5.5), Inches(1.2), BG_CARD, ACCENT4, Pt(2))
add_textbox(slide, Inches(0.7), Inches(4.1), Inches(5), Inches(0.6), "Disruptor 快", font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0.7), Inches(4.5), Inches(5), Inches(0.6), "11.6x", font_size=42, color=ACCENT4, bold=True, alignment=PP_ALIGN.CENTER)

# 3P1C Results
add_textbox(slide, Inches(6.8), Inches(1.3), Inches(5.5), Inches(0.4), "多生产者-单消费者 (3P-1C)", font_size=18, color=ACCENT2, bold=True)
bar_data_3p1c = [("Disruptor", 71, ACCENT4), ("ArrayBlockingQueue", 122, ACCENT3)]
for i, (name, ms, color) in enumerate(bar_data_3p1c):
    y = Inches(1.9) + i * Inches(1.0)
    add_textbox(slide, Inches(6.8), y, Inches(2.5), Inches(0.35), name, font_size=13, color=WHITE, bold=True)
    bar_width = max(Inches(0.3), Inches(ms / 122.0 * 3.5))
    add_rect(slide, Inches(9.5), y + Inches(0.05), bar_width, Inches(0.35), color)
    add_textbox(slide, Inches(9.5) + bar_width + Inches(0.1), y, Inches(1.5), Inches(0.35), f"{ms} ms", font_size=14, color=color, bold=True)

# Best practices
add_textbox(slide, Inches(6.8), Inches(4.0), Inches(5.5), Inches(0.3), "常见陷阱", font_size=17, color=ACCENT3, bold=True)
traps = [
    ("保存 Event 引用", "Event 对象会被复用覆盖! 需复制数据"),
    ("忘记 publish", "必须在 finally 中 publish, 否则 sequence 断裂"),
    ("消费者处理太慢", "慢消费者阻塞生产者, 考虑 WorkerPool 并行"),
    ("未设异常处理", "未捕获异常导致消费者停止, 需设 ExceptionHandler"),
]
for i, (trap, solution) in enumerate(traps):
    y = Inches(4.4) + i * Inches(0.7)
    add_textbox(slide, Inches(6.8), y, Inches(0.3), Inches(0.3), "✗", font_size=14, color=ACCENT3, bold=True)
    add_textbox(slide, Inches(7.15), y, Inches(1.8), Inches(0.3), trap, font_size=12, color=WHITE, bold=True)
    add_textbox(slide, Inches(7.15), y + Inches(0.25), Inches(5), Inches(0.3), solution, font_size=11, color=LIGHT_GRAY)


# ╔═══════════════════════════════════════════════════════════════╗
# ║           CHAPTER 8: 手写简易 Disruptor                       ║
# ╚═══════════════════════════════════════════════════════════════╝

slide = new_slide()
add_section_title(slide, "08  手写简易 Disruptor")

slide = new_slide()
add_chapter_header(slide, 8, "MiniDisruptor 核心组件", "从零实现核心功能, 彻底理解原理")

components = [
    ("MiniSequence", "带缓存行填充的序号\n前后各7个long padding\nVarHandle实现CAS", ACCENT),
    ("MiniRingBuffer", "预分配环形数组\n位运算取模\n追尾保护", ACCENT4),
    ("MiniBatchProcessor", "批量事件处理器\nwaitFor → 批量消费\n更新消费进度", ORANGE),
    ("MiniWaitStrategies", "3种等待策略实现\nYielding/Sleeping\n/BusySpin", ACCENT2),
    ("MiniDisruptor", "门面整合类\n创建/注册/启动\n/关闭", YELLOW),
    ("性能测试结果", "100万事件测试\n~9.4M ops/s\n接近官方性能", ACCENT3),
]
for i, (name, desc, color) in enumerate(components):
    col = i % 3
    row = i // 3
    x = Inches(0.4) + col * Inches(4.2)
    y = Inches(1.4) + row * Inches(2.7)
    card = add_rect(slide, x, y, Inches(4.0), Inches(2.4), BG_CARD, color, Pt(1.5))
    add_textbox(slide, x + Inches(0.2), y + Inches(0.12), Inches(3.6), Inches(0.35), name, font_size=16, color=color, bold=True, font_name="Consolas")
    add_textbox(slide, x + Inches(0.2), y + Inches(0.55), Inches(3.6), Inches(1.6), desc, font_size=13, color=LIGHT_GRAY)

# --- 8-2: MiniSequence code ---
slide = new_slide()
add_chapter_header(slide, 8, "核心代码: MiniSequence + MiniRingBuffer")

add_code_block(slide, Inches(0.4), Inches(1.3), Inches(6.2), Inches(5.8),
"""// MiniSequence - 带缓存行填充
public class MiniSequence {
    // 前置填充 (7 × 8 = 56 bytes)
    private long p1,p2,p3,p4,p5,p6,p7;

    private volatile long value; // 实际值

    // 后置填充 (7 × 8 = 56 bytes)
    private long p8,p9,p10,p11,p12,p13,p14;

    // VarHandle CAS (Java 9+ 推荐方式)
    private static final VarHandle VALUE_HANDLE;
    static {
        VALUE_HANDLE = MethodHandles.lookup()
            .findVarHandle(MiniSequence.class,
                           "value", long.class);
    }

    public long get() { return value; }
    public void set(long v) { this.value = v; }
    public boolean compareAndSet(
            long expected, long update) {
        return VALUE_HANDLE.compareAndSet(
            this, expected, update);
    }
}""", font_size=11)

add_code_block(slide, Inches(6.8), Inches(1.3), Inches(6.2), Inches(5.8),
"""// MiniRingBuffer - 环形缓冲区核心
public class MiniRingBuffer<E> {
    private final Object[] entries; // 预分配
    private final int indexMask;    // bufferSize-1
    private final MiniSequence cursor =
        new MiniSequence(-1);

    public MiniRingBuffer(MiniEventFactory<E> f,
                          int bufferSize) {
        this.indexMask = bufferSize - 1;
        this.entries = new Object[bufferSize];
        // 预填充所有 Event 对象
        for (int i = 0; i < bufferSize; i++)
            entries[i] = f.newInstance();
    }

    public E get(long sequence) {
        // 位运算取模!
        return (E) entries[
            (int)(sequence & indexMask)];
    }

    public long next() {
        long nextSeq = nextValue + 1;
        long wrapPoint = nextSeq - bufferSize;
        if (wrapPoint > cachedGatingValue) {
            // 等待消费者跟上... (追尾保护)
        }
        nextValue = nextSeq;
        return nextSeq;
    }

    public void publish(long seq) {
        cursor.set(seq); // volatile写,通知消费者
    }
}""", font_size=10)


# ╔═══════════════════════════════════════════════════════════════╗
# ║           REAL WORLD: 订单系统实战                             ║
# ╚═══════════════════════════════════════════════════════════════╝

slide = new_slide()
add_chapter_header(slide, 8, "实战案例: 高性能订单处理系统", "Demo08: 菱形 + 串行混合管道")

add_code_block(slide, Inches(0.4), Inches(1.3), Inches(7), Inches(2.2),
"""处理流程 (菱形 + 串行混合):

  [订单接入] → [参数校验] → [风控检查] ─┐
                           → [库存锁定] ─┤
                                         ├→ [价格计算] → [订单确认] → [清理]

DSL 代码:
  disruptor.handleEventsWith(validationHandler)
           .then(riskCheckHandler, inventoryHandler)  // 并行
           .then(pricingHandler)                      // 等待两者完成
           .then(confirmHandler)
           .then(cleanHandler);""", font_size=12)

# Handler cards
handlers = [
    ("1. 参数校验", "校验订单参数有效性", ACCENT4),
    ("2. 风控检查", "金额>5万 → 高风险", ORANGE),
    ("3. 库存锁定", "数量>4 → 库存不足", ORANGE),
    ("4. 价格计算", "总额+折扣计算", ACCENT),
    ("5. 订单确认", "最终确认+延迟统计", ACCENT2),
    ("6. 事件清理", "释放大对象引用", LIGHT_GRAY),
]
for i, (name, desc, color) in enumerate(handlers):
    col = i % 3
    row = i // 3
    x = Inches(0.4) + col * Inches(4.2)
    y = Inches(3.8) + row * Inches(1.5)
    card = add_rect(slide, x, y, Inches(4.0), Inches(1.2), BG_CARD, color, Pt(1))
    add_textbox(slide, x + Inches(0.15), y + Inches(0.1), Inches(3.7), Inches(0.3), name, font_size=14, color=color, bold=True)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.45), Inches(3.7), Inches(0.6), desc, font_size=12, color=LIGHT_GRAY)

# Right side: application scenarios
add_textbox(slide, Inches(7.8), Inches(1.3), Inches(5), Inches(0.3), "Disruptor 应用场景", font_size=16, color=YELLOW, bold=True)
scenarios = [
    "LMAX Exchange - 外汇交易核心引擎",
    "Log4j2 - AsyncLogger 底层实现",
    "消息中间件 - MQ 内部高性能缓冲",
    "实时数据处理 - 高频数据管道",
    "游戏服务器 - 高并发事件处理",
]
for i, s in enumerate(scenarios):
    y = Inches(1.7) + i * Inches(0.4)
    add_textbox(slide, Inches(7.8), y, Inches(0.3), Inches(0.3), "▸", font_size=12, color=YELLOW)
    add_textbox(slide, Inches(8.15), y, Inches(4.5), Inches(0.35), s, font_size=13, color=WHITE)


# ╔═══════════════════════════════════════════════════════════════╗
# ║           WHY FAST: 总结一页                                  ║
# ╚═══════════════════════════════════════════════════════════════╝

slide = new_slide()
add_chapter_header(slide, 0, "为什么 Disruptor 这么快？", "六大核心优化手段")

optimizations = [
    ("预分配 · 零GC", "Event 对象预创建, 循环复用\n发布时只覆盖, 不 new 对象\n零内存分配, 老年代长期存活", ACCENT4, "01"),
    ("数组 · 缓存友好", "连续内存布局\nCPU 预取高效\n顺序访问命中率高", ACCENT, "02"),
    ("CacheLine Padding", "前后各 7 个 long 填充\nvalue 独占一个缓存行\n消除伪共享", YELLOW, "03"),
    ("CAS 无锁", "SingleProducer: 零竞争\nMultiProducer: CAS\nvolatile 保证可见性", ORANGE, "04"),
    ("可插拔 WaitStrategy", "BusySpin → 纳秒延迟\nYielding → 亚微秒\nBlocking → 节省 CPU", ACCENT2, "05"),
    ("批量消费 Batching", "一次 waitFor 返回多个序号\n减少系统调用和上下文切换\nendOfBatch 批量提交", ACCENT3, "06"),
]
for i, (title, desc, color, num) in enumerate(optimizations):
    col = i % 3
    row = i // 3
    x = Inches(0.3) + col * Inches(4.3)
    y = Inches(1.3) + row * Inches(3.0)
    card = add_rect(slide, x, y, Inches(4.1), Inches(2.7), BG_CARD, color, Pt(2))
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.15), Inches(0.5), Inches(0.5))
    circle.fill.solid(); circle.fill.fore_color.rgb = color; circle.line.fill.background()
    circle.text_frame.paragraphs[0].text = num
    circle.text_frame.paragraphs[0].font.size = Pt(16); circle.text_frame.paragraphs[0].font.color.rgb = BG_DARK
    circle.text_frame.paragraphs[0].font.bold = True; circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_textbox(slide, x + Inches(0.75), y + Inches(0.2), Inches(3.1), Inches(0.4), title, font_size=16, color=color, bold=True)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.75), Inches(3.7), Inches(1.7), desc, font_size=13, color=LIGHT_GRAY)


# ╔═══════════════════════════════════════════════════════════════╗
# ║           FINAL: Thank you / Q&A                             ║
# ╚═══════════════════════════════════════════════════════════════╝

slide = new_slide()
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT)
add_rect(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), ACCENT2)

add_textbox(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(1),
            "Thank You", font_size=56, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(5.5), Inches(3.0), Inches(2.3), Inches(0.04), ACCENT)
add_textbox(slide, Inches(1), Inches(3.3), Inches(11.3), Inches(0.6),
            "Q & A", font_size=36, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

# Resources
resources = [
    "GitHub 教程: github.com/malinghan/madisruptor",
    "LMAX Disruptor: github.com/LMAX-Exchange/disruptor",
    "技术论文: lmax-exchange.github.io/disruptor/disruptor.html",
    "Mechanical Sympathy Blog: mechanical-sympathy.blogspot.com",
]
for i, r in enumerate(resources):
    add_textbox(slide, Inches(3.5), Inches(4.3) + i * Inches(0.45), Inches(6.3), Inches(0.4),
                r, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER, font_name="Consolas")


# ═══════════════════════════════════════════════════════════════
#  SAVE
# ═══════════════════════════════════════════════════════════════

output_path = "/home/user/webapp/Disruptor原理详解.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
